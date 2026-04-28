"""
Index OCL documents into Airtable Knowledge_Base.

Reads PDFs, Word, PowerPoint, and text files — extracts text,
uses Claude to generate summary + key insights + objection handlers,
then upserts into Knowledge_Base so agents can search before LLM calls.

Usage:
    python3 scripts/index_knowledge_base.py --dry-run   # preview only
    python3 scripts/index_knowledge_base.py             # index everything
    python3 scripts/index_knowledge_base.py --path /path/to/file.pdf  # single file
    python3 scripts/index_knowledge_base.py --reindex   # force reindex existing

The agent can also trigger this from Airtable:
  Drop a file into Knowledge_Base.Source_File → webhook → this script
  (set Status=pending_extraction, script picks it up)
"""

import argparse
import logging
import os
import sys
import time
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv
load_dotenv(Path(__file__).parent.parent / ".env")

logging.basicConfig(level=logging.INFO, format="%(levelname)s  %(message)s")
log = logging.getLogger("index_kb")

sys.path.insert(0, str(Path(__file__).parent.parent))

AIRTABLE_TOKEN  = os.getenv("AIRTABLE_TOKEN")
AIRTABLE_BASE   = os.getenv("AIRTABLE_BASE_ID", "appwlGCHBaRef70gO")
ANTHROPIC_KEY   = os.getenv("ANTHROPIC_API_KEY")
KB_TABLE        = "Knowledge_Base"

# Map filename patterns → doc_type
DOC_TYPE_RULES = [
    (["5p", "5ps", "five_p"],            "skill_packet"),
    (["kcpov", "pov", "email_template"], "skill_packet"),
    (["persona", "personas", "buyer"],   "persona_guide"),
    (["case_study", "case study"],       "case_study"),
    (["transcript", "call_recording"],   "call_transcript"),
    (["gartner", "chro", "hr_trend",
      "hrci", "state_of_hr"],            "research_report"),
    (["storybrand", "story_brand"],      "framework"),
    (["sales_story", "our_why",
      "brochure", "why oncourse"],       "playbook"),
    (["testimonial"],                    "case_study"),
    (["objection", "objections"],        "objection_handler"),
    (["competitor", "bai", "skillsoft"], "competitor_intel"),
]

ICP_LANE_RULES = [
    (["compliance", "bsa", "aml",
      "regulatory", "fair_lending"],     "compliance"),
    (["chro", "hr", "human_resources",
      "people"],                         "hr"),
    (["l&d", "learning", "training",
      "l_and_d"],                        "l_and_d"),
]


# ── Text extraction ────────────────────────────────────────────────────

def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        elif ext in (".docx", ".doc"):
            return _extract_docx(path)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(path)
        elif ext in (".txt", ".md"):
            return path.read_text(encoding="utf-8", errors="ignore")
        else:
            log.warning(f"Unsupported type: {ext}")
            return ""
    except Exception as e:
        log.error(f"Extraction failed for {path.name}: {e}")
        return ""


def _extract_pdf(path: Path) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        if len(lines) > 1:
            slides.append("\n".join(lines))
    return "\n\n".join(slides)


# ── Classification helpers ─────────────────────────────────────────────

def classify_doc_type(name: str) -> str:
    name_lower = name.lower().replace(" ", "_").replace("-", "_")
    for keywords, doc_type in DOC_TYPE_RULES:
        if any(kw in name_lower for kw in keywords):
            return doc_type
    return "other"


def classify_icp_lane(name: str, text: str) -> str:
    combined = (name + " " + text[:500]).lower()
    for keywords, lane in ICP_LANE_RULES:
        if any(kw in combined for kw in keywords):
            return lane
    return "all"


def classify_relevant_to(doc_type: str) -> str:
    mapping = {
        "skill_packet":      "FORGE",
        "framework":         "FORGE",
        "email_example":     "FORGE",
        "persona_guide":     "RECON",
        "research_report":   "RECON",
        "case_study":        "RECON",
        "call_transcript":   "CORTEX",
        "objection_handler": "CORTEX",
        "competitor_intel":  "RECON",
        "playbook":          "all_agents",
        "win_loss":          "CORTEX",
    }
    return mapping.get(doc_type, "all_agents")


# ── Claude analysis ────────────────────────────────────────────────────

def analyze_with_claude(title: str, doc_type: str, text: str) -> dict:
    """Call Claude to generate summary, insights, and objection handlers."""
    if not ANTHROPIC_KEY:
        log.warning("ANTHROPIC_API_KEY not set — skipping Claude analysis")
        return {
            "Summary":           "(Claude analysis pending — add ANTHROPIC_API_KEY)",
            "Key_Insights":      "",
            "Objection_Handles": "",
        }

    import anthropic
    client = anthropic.Anthropic(api_key=ANTHROPIC_KEY)

    # Truncate to ~12k tokens for context
    truncated = text[:50_000] if len(text) > 50_000 else text

    prompt = f"""You are analyzing a sales enablement document for ATLAS, an AI outbound agent
targeting US community banks and credit unions for compliance and HR training sales.

Document: "{title}"
Type: {doc_type}

FULL TEXT:
{truncated}

Return a JSON object with exactly these keys:
{{
  "summary": "3-5 sentence summary of what this document is and why it matters for outbound sales",
  "key_insights": "10-15 bullet points (one per line, start each with •) of the most actionable insights ATLAS agents should know",
  "objection_handles": "If this doc contains objections or objection responses, extract them as: OBJECTION: ... | RESPONSE: ... (one per line). Otherwise return empty string.",
  "tags": "comma-separated tags relevant to this doc (e.g. bsa, renewal_window, fair_lending, kcpov, 5p)"
}}

Return ONLY valid JSON. No markdown, no explanation."""

    try:
        msg = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=2000,
            messages=[{"role": "user", "content": prompt}],
        )
        import json
        raw = msg.content[0].text.strip()
        data = json.loads(raw)
        return {
            "Summary":           data.get("summary", ""),
            "Key_Insights":      data.get("key_insights", ""),
            "Objection_Handles": data.get("objection_handles", ""),
            "Tags":              data.get("tags", ""),
        }
    except Exception as e:
        log.error(f"Claude analysis failed: {e}")
        return {
            "Summary":           f"(Analysis failed: {e})",
            "Key_Insights":      "",
            "Objection_Handles": "",
        }


# ── Airtable upsert ────────────────────────────────────────────────────

def upsert_to_kb(record: dict) -> bool:
    import requests
    from datetime import datetime, timezone
    url = f"https://api.airtable.com/v0/{AIRTABLE_BASE}/{KB_TABLE}"
    headers = {"Authorization": f"Bearer {AIRTABLE_TOKEN}", "Content-Type": "application/json"}

    payload = {
        "records": [{"fields": record}],
        "performUpsert": {"fieldsToMergeOn": ["Title"]},
    }
    resp = requests.patch(url, json=payload, headers=headers, timeout=30)
    if not resp.ok:
        log.error(f"Airtable upsert failed: {resp.status_code} {resp.text[:200]}")
    return resp.ok


# ── Main indexer ───────────────────────────────────────────────────────

def index_file(path: Path, dry_run: bool = False) -> bool:
    log.info(f"  Processing: {path.name}")

    # Extract text
    text = extract_text(path)
    if not text.strip():
        log.warning(f"  No text extracted from {path.name} — skipping")
        return False

    token_estimate = len(text.split()) * 1.3

    # Classify
    doc_type    = classify_doc_type(path.name)
    icp_lane    = classify_icp_lane(path.name, text)
    relevant_to = classify_relevant_to(doc_type)

    log.info(f"  type={doc_type}, lane={icp_lane}, ~{int(token_estimate)} tokens")

    if dry_run:
        print(f"    [DRY RUN] Would index: {path.name}")
        print(f"    type={doc_type}, lane={icp_lane}, tokens~{int(token_estimate)}")
        print(f"    text preview: {text[:200].strip()[:100]}...")
        return True

    # Claude analysis
    analysis = analyze_with_claude(path.name, doc_type, text)

    record = {
        "Title":             path.stem,
        "Doc_Type":          doc_type,
        "Source_Path":       str(path),
        "Extracted_Text":    text[:100_000],  # Airtable longtext limit
        "Summary":           analysis.get("Summary", ""),
        "Key_Insights":      analysis.get("Key_Insights", ""),
        "Objection_Handles": analysis.get("Objection_Handles", ""),
        "Tags":              analysis.get("Tags", ""),
        "ICP_Lane":          icp_lane,
        "Relevant_To":       relevant_to,
        "Status":            "indexed",
        "Token_Count":       int(token_estimate),
        "Active":            True,
        "Version":           1,
    }

    success = upsert_to_kb(record)
    if success:
        log.info(f"  ✓ Indexed: {path.name}")
    return success


def process_pending_attachments():
    """
    Check Knowledge_Base for records with Status=pending_extraction.
    Called by Pipedream webhook when user drops a file in Airtable.
    """
    import requests
    url = f"https://api.airtable.com/v0/{AIRTABLE_BASE}/{KB_TABLE}"
    headers = {"Authorization": f"Bearer {AIRTABLE_TOKEN}"}
    params = {"filterByFormula": '{Status}="pending_extraction"'}
    resp = requests.get(url, params=params, headers=headers, timeout=15)
    if not resp.ok:
        return

    records = resp.json().get("records", [])
    log.info(f"Found {len(records)} pending attachments")

    for rec in records:
        rid    = rec["id"]
        fields = rec.get("fields", {})
        files  = fields.get("Source_File", [])
        title  = fields.get("Title", "unknown")

        for attachment in files:
            file_url = attachment.get("url")
            filename  = attachment.get("filename", "file")
            if not file_url:
                continue

            # Download attachment
            import tempfile
            file_resp = requests.get(file_url, timeout=30)
            if not file_resp.ok:
                continue

            suffix = Path(filename).suffix
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp.write(file_resp.content)
                tmp_path = Path(tmp.name)

            # Mark as extracting
            requests.patch(
                f"{url}/{rid}",
                json={"fields": {"Status": "extracting"}},
                headers={**headers, "Content-Type": "application/json"},
                timeout=15,
            )

            text       = extract_text(tmp_path)
            doc_type   = classify_doc_type(filename)
            icp_lane   = classify_icp_lane(filename, text)
            analysis   = analyze_with_claude(filename, doc_type, text)

            update = {
                "Source_Path":       filename,
                "Extracted_Text":    text[:100_000],
                "Summary":           analysis.get("Summary", ""),
                "Key_Insights":      analysis.get("Key_Insights", ""),
                "Objection_Handles": analysis.get("Objection_Handles", ""),
                "Tags":              analysis.get("Tags", ""),
                "Doc_Type":          doc_type,
                "ICP_Lane":          icp_lane,
                "Relevant_To":       classify_relevant_to(doc_type),
                "Status":            "indexed",
                "Token_Count":       int(len(text.split()) * 1.3),
                "Active":            True,
            }
            requests.patch(
                f"{url}/{rid}",
                json={"fields": update},
                headers={**headers, "Content-Type": "application/json"},
                timeout=30,
            )
            tmp_path.unlink(missing_ok=True)
            log.info(f"  ✓ Indexed attachment: {filename}")


def main():
    parser = argparse.ArgumentParser(description="Index documents into Airtable Knowledge_Base")
    parser.add_argument("--dry-run",  action="store_true", help="Preview without writing")
    parser.add_argument("--path",     type=str,            help="Single file path to index")
    parser.add_argument("--dir",      type=str,            help="Directory to scan (default: agent-brain/raw/atlas)")
    parser.add_argument("--pending",  action="store_true", help="Process pending Airtable attachments only")
    parser.add_argument("--reindex",  action="store_true", help="Re-index already-indexed docs")
    args = parser.parse_args()

    if args.pending:
        process_pending_attachments()
        return

    # Default scan directory
    scan_dir = Path(args.dir) if args.dir else Path("/Users/anmolsam/projects/agent-brain/raw/atlas")

    if args.path:
        paths = [Path(args.path)]
    else:
        paths = sorted([
            p for p in scan_dir.rglob("*")
            if p.is_file()
            and p.suffix.lower() in (".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md")
            and not p.name.startswith(".")
            and p.parent.name != "_large_pdfs"
        ])

    print(f"\n── Indexing {len(paths)} documents ─────────────────────────────")
    ok = failed = 0

    for path in paths:
        success = index_file(path, dry_run=args.dry_run)
        if success:
            ok += 1
        else:
            failed += 1
        time.sleep(0.5)   # rate limit Airtable + Claude

    print(f"\n  Indexed: {ok}   Failed: {failed}")
    if args.dry_run:
        print("  (Dry run — run without --dry-run to write to Airtable)")
    print()


if __name__ == "__main__":
    main()
